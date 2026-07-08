import pptx

prs = pptx.Presentation(r'E:\sEMG_B_Project\docs\sEMG肌肉疲劳监测_操作指南.pptx')
print(f'Total slides: {len(prs.slides)}')
print(f'Size: {prs.slide_width/914400:.2f}" x {prs.slide_height/914400:.2f}"')
print()

for i, slide in enumerate(prs.slides):
    shapes = list(slide.shapes)
    text_shapes = [s for s in shapes if s.has_text_frame and s.text_frame.text.strip()]
    pics = [s for s in shapes if s.shape_type == 13]
    tables = [s for s in shapes if s.has_table]

    # 收集主要文本
    texts = []
    for s in text_shapes[:3]:
        t = s.text_frame.text.strip()[:50]
        if t:
            texts.append(t)

    # 收集图片位置
    pic_info = []
    for p in pics:
        left = p.left / 914400
        top = p.top / 914400
        width = p.width / 914400
        height = p.height / 914400
        pic_info.append(f'img({left:.1f},{top:.1f} {width:.1f}x{height:.1f})')

    print(f'Slide {i}: {len(shapes)} shapes, {len(text_shapes)} texts, {len(pics)} pics, {len(tables)} tables')
    for t in texts:
        print(f'  Text: {t}')
    for p in pic_info:
        print(f'  {p}')
    if tables:
        print(f'  Table: {len(tables[0].table.rows)} rows x {len(tables[0].table.columns)} cols')
    print()
