import pptx

prs = pptx.Presentation(r'E:\sEMG_B_Project\docs\sEMG肌肉疲劳监测_操作指南.pptx')

for i, slide in enumerate(prs.slides):
    print(f'\n=== Slide {i} ===')
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            left = shape.left / 914400
            top = shape.top / 914400
            w = shape.width / 914400
            h = shape.height / 914400
            text = shape.text_frame.text.strip()[:80]
            print(f'  Text[{left:.1f},{top:.1f} {w:.1f}x{h:.1f}]: {text}')
        elif shape.shape_type == 13:
            left = shape.left / 914400
            top = shape.top / 914400
            w = shape.width / 914400
            h = shape.height / 914400
            print(f'  Img[{left:.1f},{top:.1f} {w:.1f}x{h:.1f}]')
        elif shape.has_table:
            print(f'  Table: {len(shape.table.rows)} rows x {len(shape.table.columns)} cols')
            for ri, row in enumerate(shape.table.rows):
                cells = [cell.text[:20] for cell in row.cells]
                print(f'    Row{ri}: {cells}')
