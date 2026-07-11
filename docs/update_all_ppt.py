#!/usr/bin/env python3
from pptx import Presentation

files = [
    'E:/sEMG_B_Project/docs/sEMG肌肉疲劳监测_小程序操作指南.pptx',
    'E:/sEMG_B_Project/docs/deck_workspace/sEMG肌肉疲劳监测_操作指南_new.pptx',
]

for f in files:
    try:
        prs = Presentation(f)
        updated = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            old = run.text
                            new = old
                            if '动态基线' in new:
                                new = new.replace('动态基线', '锚点公式')
                            if '收缩起始MDF' in new:
                                new = new.replace('收缩起始MDF', 'Active峰值MDF')
                            if '收缩起始-当前' in new:
                                new = new.replace('收缩起始-当前', 'Active峰值MDF-当前MDF')
                            if '收缩起始MDF×100%' in new:
                                new = new.replace('收缩起始MDF×100%', '(Active峰值MDF-用力末MDF)×100%')
                            if '收缩起始' in new and 'MDF' in new:
                                new = new.replace('收缩起始', 'Active峰值')
                            if new != old:
                                run.text = new
                                updated = True
        if updated:
            prs.save(f)
            print(f'Updated: {f}')
        else:
            print(f'No changes needed: {f}')
    except Exception as e:
        print(f'Error processing {f}: {e}')