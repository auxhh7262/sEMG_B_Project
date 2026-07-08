import pptx
from pptx.util import Pt, Emu

prs = pptx.Presentation(r'E:\sEMG_B_Project\docs\sEMG肌肉疲劳监测_操作指南_backup.pptx')

# 分析几个关键页面的详细样式
for slide_idx in [0, 1, 2, 3, 4, 12, 17]:
    slide = prs.slides[slide_idx]
    print(f'\n=== Slide {slide_idx} ===')
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        text = tf.text.strip()[:60]
        if not text:
            continue
        
        left = shape.left / 914400
        top = shape.top / 914400
        width = shape.width / 914400
        height = shape.height / 914400
        
        # 获取第一个段落的样式
        para = tf.paragraphs[0] if tf.paragraphs else None
        if para and para.runs:
            run = para.runs[0]
            font = run.font
            font_size = font.size.pt if font.size else 'N/A'
            font_name = font.name if font.name else 'N/A'
            bold = font.bold
            color = str(font.color.rgb) if font.color and font.color.rgb else 'N/A'
            align = para.alignment
        else:
            font_size = 'N/A'
            font_name = 'N/A'
            bold = 'N/A'
            color = 'N/A'
            align = 'N/A'
        
        print(f'  [{left:.1f},{top:.1f} {width:.1f}x{height:.1f}] "{text}"')
        print(f'    font={font_name} {font_size}pt bold={bold} color={color} align={align}')
